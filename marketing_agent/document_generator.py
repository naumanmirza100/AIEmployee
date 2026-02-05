"""
Document Generator - Converts marketing documents to various file formats
Supports: PDF, PPTX (DOCX removed)
"""

import io
import re
from typing import Optional, List, Tuple
from django.http import HttpResponse
from marketing_agent.models import MarketingDocument

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import black, HexColor, white
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak,
        Table, TableStyle,
    )
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

# Optional: charts (reportlab.graphics)
try:
    from reportlab.graphics.shapes import Drawing
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    from reportlab.graphics.charts.piecharts import Pie
    from reportlab.graphics.charts.legends import Legend
    CHARTS_AVAILABLE = True
except ImportError:
    CHARTS_AVAILABLE = False

# PDF layout constants
PDF_HEADER_HEIGHT = 0.5 * inch
PDF_FOOTER_HEIGHT = 0.45 * inch
PDF_MARGIN = 0.75 * inch


def _bold_markdown(text: str) -> str:
    """Convert **text** to <b>text</b> (all occurrences) for ReportLab Paragraph."""
    if not text:
        return text
    return re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)


def _parse_markdown_table(lines: List[str], start: int) -> Tuple[Optional[List[List[str]]], int]:
    """Parse a markdown table (header, optional separator, body rows). Returns (table_data or None, next_line_index)."""
    rows = []
    i = start
    sep_seen = False
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()
        if not stripped:
            break
        # Separator row: |---| or |:---:|---|
        if re.match(r'^[\|\s\-:]+\s*$', stripped) and re.search(r'\-+', stripped):
            if rows:
                sep_seen = True
            i += 1
            continue
        # Table row: | cell | cell |
        if '|' in stripped:
            cells = [c.strip() for c in stripped.split('|')]
            cells = [c for c in cells if c is not None]
            if not cells or all(c == '' for c in cells):
                i += 1
                continue
            rows.append(cells)
            i += 1
        else:
            if sep_seen or rows:
                break
            i += 1
            break
    if not rows:
        return None, start
    return rows, i


def _parse_chart_block(lines: List[str], start: int) -> Tuple[Optional[dict], int]:
    """Parse [CHART] ... [/CHART] or ```chart ... ``` block. Returns (dict with type, title, labels, values) or (None, start)."""
    if start >= len(lines):
        return None, start
    first = lines[start].strip().upper()
    if first in ('[CHART]', '```CHART'):
        end_marker = '[/CHART]' if first == '[CHART]' else '```'
        i = start + 1
        block = {}
        while i < len(lines):
            line = lines[i]
            stripped = line.strip()
            if stripped.upper() == end_marker or stripped == '```':
                i += 1
                break
            if ':' in stripped:
                key, _, val = stripped.partition(':')
                key = key.strip().lower()
                val = val.strip()
                if key == 'type':
                    block['type'] = val.lower()
                elif key == 'title':
                    block['title'] = val
                elif key == 'labels':
                    block['labels'] = [x.strip() for x in val.split(',') if x.strip()]
                elif key == 'values':
                    block['values'] = []
                    for x in val.split(','):
                        x = x.strip()
                        if x:
                            try:
                                block['values'].append(float(x))
                            except ValueError:
                                block['values'].append(0.0)
            i += 1
        if block.get('type') and block.get('labels') and block.get('values') and len(block['labels']) == len(block['values']):
            return block, i
        return None, start
    return None, start


def _build_chart_drawing(chart_spec: dict, styles) -> Optional["Drawing"]:
    """Build a ReportLab Drawing (bar or pie chart). Returns None if charts not available or invalid."""
    if not CHARTS_AVAILABLE:
        return None
    chart_type = (chart_spec.get('type') or 'bar').lower()
    if chart_type == 'line':
        chart_type = 'bar'  # Line not supported; use bar for trend data
    title = chart_spec.get('title', '')
    labels = chart_spec.get('labels', [])
    values = chart_spec.get('values', [])
    if not labels or not values:
        return None
    width, height = 400, 220
    drawing = Drawing(width, height)
    if chart_type == 'pie':
        pie = Pie()
        pie.x = 120
        pie.y = 30
        pie.width = 160
        pie.height = 160
        pie.data = list(values)
        pie.labels = list(labels)
        pie.slices.strokeWidth = 0.5
        colors = [HexColor('#4a5568'), HexColor('#718096'), HexColor('#a0aec0'), HexColor('#e2e8f0'), HexColor('#2d3748'), HexColor('#63b3ed')]
        for i in range(len(pie.slices)):
            pie.slices[i].fillColor = colors[i % len(colors)]
        drawing.add(pie)
    else:
        # Bar chart
        bc = VerticalBarChart()
        bc.x = 50
        bc.y = 30
        bc.width = 340
        bc.height = 160
        bc.data = [tuple(values)]
        bc.categoryAxis.categoryNames = labels
        bc.categoryAxis.labels.angle = 25
        bc.categoryAxis.labels.fontSize = 8
        bc.valueAxis.valueMin = 0
        bc.valueAxis.valueMax = max(values) * 1.15 if values else 100
        bc.bars[0].fillColor = HexColor('#4a5568')
        bc.barSpacing = 4
        bc.groupSpacing = 12
        drawing.add(bc)
    return drawing


class DocumentGenerator:
    """Generate documents in various formats"""

    @staticmethod
    def generate_pdf(document: MarketingDocument) -> HttpResponse:
        """Generate PDF file with proper header, footer, tables, and headings."""
        if not PDF_AVAILABLE:
            raise ImportError("reportlab is not installed. Run: pip install reportlab")

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=letter,
            leftMargin=PDF_MARGIN,
            rightMargin=PDF_MARGIN,
            topMargin=PDF_MARGIN + PDF_HEADER_HEIGHT,
            bottomMargin=PDF_MARGIN + PDF_FOOTER_HEIGHT,
        )

        elements = []
        styles = DocumentGenerator._pdf_styles()

        # ---- Document title (first page only, in body)
        title_style = styles['CustomTitle']
        elements.append(Paragraph(document.title, title_style))
        elements.append(Spacer(1, 0.15 * inch))

        # ---- Metadata as a small table
        meta_data = [
            ['Document Type', document.get_document_type_display()],
            ['Status', document.get_status_display()],
            ['Created', document.created_at.strftime('%Y-%m-%d %H:%M')],
        ]
        if document.campaign:
            meta_data.append(['Campaign', document.campaign.name])
        meta_table = Table(meta_data, colWidths=[1.2 * inch, 4.3 * inch])
        meta_table.setStyle(TableStyle([
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('TEXTCOLOR', (0, 0), (0, -1), HexColor('#555')),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
            ('TOPPADDING', (0, 0), (-1, -1), 4),
            ('LINEBELOW', (0, 0), (-1, -1), 0.5, HexColor('#e0e0e0')),
            ('LINEABOVE', (0, 0), (-1, 0), 0.5, HexColor('#e0e0e0')),
        ]))
        elements.append(meta_table)
        elements.append(Spacer(1, 0.35 * inch))

        # ---- Body content (headings, paragraphs, lists, tables)
        content = document.content or ''
        DocumentGenerator._add_content_to_pdf(elements, content, styles)

        # ---- Header/footer callbacks
        doc_title = document.title
        doc_date = document.created_at.strftime('%B %d, %Y')

        def draw_header(canvas, doc):
            canvas.saveState()
            canvas.setFont('Helvetica-Bold', 10)
            canvas.setFillColor(HexColor('#333'))
            canvas.drawString(PDF_MARGIN, letter[1] - PDF_MARGIN - 14, doc_title)
            canvas.setStrokeColor(HexColor('#cccccc'))
            canvas.setLineWidth(0.5)
            canvas.line(PDF_MARGIN, letter[1] - PDF_MARGIN - 18, letter[0] - PDF_MARGIN, letter[1] - PDF_MARGIN - 18)
            canvas.restoreState()

        def draw_footer(canvas, doc):
            canvas.saveState()
            canvas.setFont('Helvetica', 8)
            canvas.setFillColor(HexColor('#666'))
            canvas.drawString(PDF_MARGIN, PDF_MARGIN - 10, doc_date)
            canvas.drawRightString(letter[0] - PDF_MARGIN, PDF_MARGIN - 10, f"Page {doc.page}")
            canvas.setStrokeColor(HexColor('#cccccc'))
            canvas.setLineWidth(0.5)
            canvas.line(PDF_MARGIN, PDF_MARGIN + 4, letter[0] - PDF_MARGIN, PDF_MARGIN + 4)
            canvas.restoreState()

        def first_page(canvas, doc):
            draw_header(canvas, doc)
            draw_footer(canvas, doc)

        def later_pages(canvas, doc):
            draw_header(canvas, doc)
            draw_footer(canvas, doc)

        doc.build(elements, onFirstPage=first_page, onLaterPages=later_pages)

        pdf = buffer.getvalue()
        buffer.close()

        response = HttpResponse(pdf, content_type='application/pdf')
        response['Content-Disposition'] = f'attachment; filename="{document.title.replace(" ", "_")}.pdf"'
        return response

    @staticmethod
    def _pdf_styles():
        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            textColor=HexColor('#1a1a1a'),
            spaceAfter=12,
            alignment=TA_LEFT,
        ))
        styles.add(ParagraphStyle(
            name='CustomH1',
            parent=styles['Heading1'],
            fontSize=14,
            textColor=HexColor('#2c2c2c'),
            spaceBefore=14,
            spaceAfter=8,
            borderPadding=(0, 0, 0, 0),
        ))
        styles.add(ParagraphStyle(
            name='CustomH2',
            parent=styles['Heading2'],
            fontSize=12,
            textColor=HexColor('#404040'),
            spaceBefore=10,
            spaceAfter=6,
        ))
        styles.add(ParagraphStyle(
            name='CustomH3',
            parent=styles['Heading3'],
            fontSize=11,
            textColor=HexColor('#555'),
            spaceBefore=8,
            spaceAfter=4,
        ))
        styles.add(ParagraphStyle(
            name='CustomNormal',
            parent=styles['Normal'],
            fontSize=10,
            leading=12,
            spaceAfter=6,
        ))
        styles.add(ParagraphStyle(
            name='CustomBoldSubhead',
            parent=styles['Normal'],
            fontName='Helvetica-Bold',
            fontSize=10,
            leading=12,
            spaceBefore=6,
            spaceAfter=4,
        ))
        return styles

    @staticmethod
    def _add_content_to_pdf(elements, content: str, styles):
        """Add formatted content: headings, paragraphs, lists, markdown tables."""
        lines = content.split('\n')
        i = 0
        while i < len(lines):
            line = lines[i]
            stripped = line.strip()
            if not stripped:
                elements.append(Spacer(1, 0.12 * inch))
                i += 1
                continue

            # Chart block: [CHART] ... [/CHART] or ```chart ... ```
            if stripped.upper() in ('[CHART]', '```CHART'):
                chart_spec, next_i = _parse_chart_block(lines, i)
                if chart_spec:
                    drawing = _build_chart_drawing(chart_spec, styles)
                    if drawing:
                        if chart_spec.get('title'):
                            elements.append(Paragraph(_bold_markdown(chart_spec['title']), styles['CustomH3']))
                            elements.append(Spacer(1, 0.08 * inch))
                        elements.append(drawing)
                        elements.append(Spacer(1, 0.2 * inch))
                    i = next_i
                    continue

            # Markdown table: must start with | or look like table
            if stripped.startswith('|') or (i + 1 < len(lines) and '|' in stripped and '|' in lines[i + 1].strip()):
                table_data, next_i = _parse_markdown_table(lines, i)
                if table_data and len(table_data) > 0:
                    # Build Table flowable with header row styled; constrain width so table fits on page
                    col_count = max(len(r) for r in table_data)
                    for row in table_data:
                        while len(row) < col_count:
                            row.append('')
                    avail_width = 6.5 * inch  # letter width minus margins
                    col_width = avail_width / max(col_count, 1)
                    col_widths = [col_width] * col_count
                    # Wrap cell content in Paragraph so long text wraps and does not overflow
                    try:
                        wrap_style = styles['CustomNormal']
                        table_data_wrapped = [[Paragraph(str(cell)[:500], wrap_style) for cell in row] for row in table_data]
                    except Exception:
                        table_data_wrapped = table_data
                    table = Table(table_data_wrapped, colWidths=col_widths, repeatRows=1)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#4a5568')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), white),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 9),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                        ('INNERGRID', (0, 0), (-1, -1), 0.5, HexColor('#e2e8f0')),
                        ('BOX', (0, 0), (-1, -1), 0.5, HexColor('#cbd5e0')),
                        ('LEFTPADDING', (0, 0), (-1, -1), 8),
                        ('RIGHTPADDING', (0, 0), (-1, -1), 8),
                        ('TOPPADDING', (0, 0), (-1, -1), 6),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8fafc')]),
                    ]))
                    elements.append(table)
                    elements.append(Spacer(1, 0.2 * inch))
                    i = next_i
                    continue

            # Headings: strip # with or without space so "#" never appears in PDF (Marketing Strategy and all docs)
            # ## = H2 (main sections only for strategy); # = H1; ###/#### = bold subhead
            hash_match = re.match(r'^(#+)\s*(.*)$', stripped)
            if hash_match and stripped[0] == '#':
                prefix, rest = hash_match.group(1), hash_match.group(2).strip()
                num_hashes = len(prefix)
                if num_hashes == 1:
                    elements.append(Paragraph(_bold_markdown(rest), styles['CustomH1']))
                    elements.append(Spacer(1, 0.08 * inch))
                elif num_hashes == 2:
                    elements.append(Paragraph(_bold_markdown(rest), styles['CustomH2']))
                    elements.append(Spacer(1, 0.06 * inch))
                else:
                    # ### or #### or more: render as bold paragraph (no H3/H4 in PDF)
                    elements.append(Paragraph(_bold_markdown(rest), styles['CustomBoldSubhead']))
                    elements.append(Spacer(1, 0.05 * inch))
                i += 1
                continue

            # Lists (bold ** in list text)
            if stripped.startswith('- ') or stripped.startswith('* '):
                para = Paragraph("• " + _bold_markdown(stripped[2:]), styles['CustomNormal'])
                elements.append(para)
                elements.append(Spacer(1, 0.04 * inch))
                i += 1
                continue
            if re.match(r'^\d+\.\s', stripped):
                para = Paragraph(_bold_markdown(stripped), styles['CustomNormal'])
                elements.append(para)
                elements.append(Spacer(1, 0.04 * inch))
                i += 1
                continue

            # Horizontal rule
            if stripped.startswith('---') or stripped.startswith('==='):
                elements.append(Spacer(1, 0.15 * inch))
                i += 1
                continue

            # Regular paragraph (bold **text**)
            para = Paragraph(_bold_markdown(stripped), styles['CustomNormal'])
            elements.append(para)
            elements.append(Spacer(1, 0.08 * inch))
            i += 1
    
    @staticmethod
    def generate_pptx(document: MarketingDocument) -> HttpResponse:
        """Generate PPTX file from marketing document (for presentations)"""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is not installed. Run: pip install python-pptx")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        content = document.content
        
        # Split content into slides
        slides_content = DocumentGenerator._parse_presentation_slides(content)
        
        if not slides_content:
            # If no slides found, create a single slide with all content
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            DocumentGenerator._add_text_to_slide(slide, document.title, is_title=True)
            DocumentGenerator._add_text_to_slide(slide, content, is_title=False)
        else:
            # Create slides from parsed content
            for i, slide_content in enumerate(slides_content):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                
                # Extract title and content
                title, body = DocumentGenerator._extract_slide_title_body(slide_content)
                
                if title:
                    DocumentGenerator._add_text_to_slide(slide, title, is_title=True)
                if body:
                    DocumentGenerator._add_text_to_slide(slide, body, is_title=False)
        
        # Create response
        response = HttpResponse(
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response['Content-Disposition'] = f'attachment; filename="{document.title.replace(" ", "_")}.pptx"'
        
        prs.save(response)
        return response
    
    @staticmethod
    def _parse_presentation_slides(content: str) -> list:
        """Parse presentation content into individual slides"""
        # Split by slide markers
        slides = re.split(r'(?:^|\n)#\s*Slide\s+\d+:|(?:^|\n)Slide\s+\d+:|(?:^|\n)##\s*Slide\s+\d+:', content, flags=re.IGNORECASE | re.MULTILINE)
        
        # Filter out empty slides
        return [slide.strip() for slide in slides if slide.strip()]
    
    @staticmethod
    def _extract_slide_title_body(slide_content: str) -> tuple:
        """Extract title and body from slide content"""
        lines = slide_content.split('\n')
        title = None
        body_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # First non-empty line is usually the title
            if not title and line and not line.startswith('-') and not line.startswith('*'):
                if len(line) < 100:  # Reasonable title length
                    title = line
                    continue
            
            body_lines.append(line)
        
        body = '\n'.join(body_lines)
        return title, body
    
    @staticmethod
    def _add_text_to_slide(slide, text: str, is_title: bool = False):
        """Add text to a slide"""
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        # Get slide dimensions
        left = Inches(0.5)
        top = Inches(1) if is_title else Inches(2)
        width = Inches(9)
        height = Inches(5) if is_title else Inches(4.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Clear default paragraph
        text_frame.clear()
        
        # Add text
        p = text_frame.paragraphs[0]
        p.text = text[:500]  # Limit text length
        p.font.size = Pt(44) if is_title else Pt(18)
        p.font.bold = is_title
        p.alignment = PP_ALIGN.LEFT
        
        # Split long text into multiple paragraphs
        if len(text) > 500:
            remaining = text[500:]
            for chunk in [remaining[i:i+200] for i in range(0, len(remaining), 200)]:
                p = text_frame.add_paragraph()
                p.text = chunk
                p.font.size = Pt(18)
                p.alignment = PP_ALIGN.LEFT
    
    @staticmethod
    def get_available_formats(document_type: str) -> list:
        """Get available download formats for a document type (PDF and optionally PPTX)."""
        formats = []
        if PDF_AVAILABLE:
            formats.append(('pdf', 'PDF'))
        if document_type == 'presentation' and PPTX_AVAILABLE:
            formats.append(('pptx', 'PowerPoint (PPTX)'))
        return formats

