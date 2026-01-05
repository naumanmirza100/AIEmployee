from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# Configuration
THEME_BLUE = RGBColor(0, 112, 192)  # Professional Blue
THEME_LIGHT_BLUE = RGBColor(220, 235, 250)
THEME_DARK_GRAY = RGBColor(50, 50, 50)
FONT_NAME = 'Calibri' # Standard academic sans-serif

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Helper Functions
    def add_title_slide(title, subtitle, presenter_name="Your Name", date="Date"):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        # Formatting Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_para = title_shape.text_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.color.rgb = THEME_BLUE
        title_para.font.name = FONT_NAME
        
        # Formatting Subtitle
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        sub_para = subtitle_shape.text_frame.paragraphs[0]
        sub_para.font.size = Pt(24)
        sub_para.font.color.rgb = THEME_DARK_GRAY
        sub_para.font.name = FONT_NAME
        
        # Add Footer info
        left = Inches(1)
        top = Inches(6.5)
        txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(0.5))
        tf = txBox.text_frame
        tf.text = f"{presenter_name} | {date}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = THEME_DARK_GRAY
        p.alignment = PP_ALIGN.RIGHT
        
        add_notes(slide, "Introduction to SLR(1) Parsing. Welcome to the presentation.")

    def add_notes(slide, text):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = text

    def add_section_header(slide, title):
        slide.shapes.title.text = title
        slide.shapes.placeholders[1].text = "" # Clear body

    def draw_flowchart(slide):
        # Draw a simple parsing flowchart
        x_start = 1
        y = 2
        w = 1.5
        h = 0.8
        gap = 0.5
        
        # Nodes
        labels = ["Start", "Read Input", "Action (Shift/Reduce)", "Accept/Reject"]
        shapes = []
        for i, label in enumerate(labels):
            x = x_start + (i * (w + gap))
            if i == 2: x += 0.5 # Gap for decision
            
            if i == 0:
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            elif i == 2:
                shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            
            shape.fill.solid()
            shape.fill.fore_color.rgb = THEME_LIGHT_BLUE
            shape.line.color.rgb = THEME_BLUE
            shape.line.width = Pt(2)
            
            tf = shape.text_frame
            tf.text = label
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(10)
            shapes.append(shape)

        # Arrows
        for i in range(len(shapes)-1):
            start = shapes[i]
            end = shapes[i+1]
            slide.shapes.add_connector(MSO_SHAPE.ELLIPSE, start.right, start.top+start.height/2, end.left, end.top+end.height/2)

    # --- SLIDE CREATION ---

    # Slide 1: Title
    add_title_slide("SLR(1) Parsing in Compiler Construction", 
                    "Bottom-Up Parsing, Examples, Conflicts, and Resolutions")

    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[5] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Introduction to SLR(1)"
    p.font.size = Pt(36)
    p.font.color.rgb = THEME_BLUE
    p.font.bold = True
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    points = [
        "What is SLR(1)?",
        "  Simple LR(1) parsing technique.",
        "  Bottom-up parser using Shift-Reduce actions.",
        "",
        "Key Components:",
        "  • Input Buffer & Stack",
        "  • Parsing Table (ACTION & GOTO)",
        "",
        "Advantages:",
        "  • Efficient (LALR subset).",
        "  • Handles Left Recursion.",
        "",
        "Limitations:",
        "  • Cannot parse all CFGs.",
        "  • Subject to Shift/Reduce conflicts."
    ]
    for pt in points:
        p = tf.add_paragraph()
        p.text = pt
        p.font.size = Pt(18)
        p.font.name = FONT_NAME
        p.level = 1 if pt.startswith("•") else 0

    draw_flowchart(slide)
    add_notes(slide, "Explain the components. Flowchart shows: Start -> Input -> Decision -> Accept.")

    # Slide 3: Steps to Construct Table
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Steps to Construct SLR(1) Parsing Table"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    steps = [
        "1. Augment the Grammar (Add start symbol S').",
        "2. Construct LR(0) Items (productions with dots).",
        "3. Build Canonical Collection of Items (DFA States).",
        "   • Use Closure(I) and Goto(I, X).",
        "4. Compute FOLLOW sets for Non-Terminals.",
        "5. Fill Parsing Table:",
        "   • ACTION[State, Terminal] = Shift/Reduce/Accept",
        "   • GOTO[State, Non-Terminal] = Next State"
    ]

    for i, step in enumerate(steps):
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3 + i*0.5), Inches(6), Inches(0.5))
        tb.text_frame.text = step
        tb.text_frame.paragraphs[0].font.size = Pt(18)
        tb.text_frame.paragraphs[0].font.name = FONT_NAME

    # Simple Diagram for DFA Concept
    dfa_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(2.5), Inches(3.5))
    dfa_box.fill.solid()
    dfa_box.fill.fore_color.rgb = THEME_LIGHT_BLUE
    dfa_box.line.color.rgb = THEME_BLUE
    
    dfa_lbl = dfa_box.text_frame
    dfa_lbl.text = "State Transition Graph\n(DFA)"
    dfa_lbl.paragraphs[0].font.size = Pt(14)
    dfa_lbl.paragraphs[0].alignment = PP_ALIGN.CENTER
    dfa_lbl.paragraphs[0].font.bold = True
    
    # Draw 3 nodes inside
    n1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), Inches(3), Inches(0.5), Inches(0.5))
    n1.fill.fore_color.rgb = RGBColor(255,255,255)
    n1.text_frame.text = "0"
    n1.text_frame.paragraphs[0].font.size = Pt(10)
    
    n2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(4.5), Inches(0.5), Inches(0.5))
    n2.fill.fore_color.rgb = RGBColor(255,255,255)
    n2.text_frame.text = "1"
    n2.text_frame.paragraphs[0].font.size = Pt(10)

    slide.shapes.add_connector(MSO_SHAPE.LINE, n1.right, n1.top, n2.left, n2.bottom)

    add_notes(slide, "Step 3 is the core: building the DFA of items.")

    # Slide 4: FIRST and FOLLOW
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "FIRST and FOLLOW Sets"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Definitions
    def_text = "FIRST(α): Set of terminals that can begin strings derived from α.\nFOLLOW(A): Set of terminals that can appear immediately to the right of A."
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tb.text_frame.text = def_text
    tb.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Table
    rows, cols = 4, 2
    left = Inches(2)
    top = Inches(2.5)
    width = Inches(6)
    height = Inches(1.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    table.cell(0,0).text = "Non-Terminal"
    table.cell(0,1).text = "FOLLOW Set Example"
    table.rows[0].height = Inches(0.4)
    
    # Content
    data = [
        ["E", "FOLLOW(E) = {+, ), $}"],
        ["T", "FOLLOW(T) = {+, *, ), $}"],
        ["F", "FOLLOW(F) = {+, *, ), $}"]
    ]
    
    for i in range(rows-1):
        table.cell(i+1, 0).text = data[i][0]
        table.cell(i+1, 1).text = data[i][1]
        table.rows[i+1].height = Inches(0.4)
    
    # Style table
    for cell in table.iter_cells():
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.border.color.rgb = RGBColor(150,150,150)
        
    table.cell(0,0).fill.fore_color.rgb = THEME_BLUE
    table.cell(0,1).fill.fore_color.rgb = THEME_BLUE
    table.cell(0,0).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    table.cell(0,1).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)

    add_notes(slide, "Explain how FIRST and FOLLOW are computed. Crucial for fill ACTION table reductions.")

    # Slide 5: Solved Example - Grammar
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Solved Example: Grammar & Sets"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Grammar
    g_text = "Augmented Grammar:\n1. E' -> E\n2. E -> E + T\n3. E -> T\n4. T -> T * F\n5. T -> F\n6. F -> ( E )\n7. F -> id"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = g_text
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.name = "Courier New" # Monospace for grammar

    # Table for Sets
    rows, cols = 4, 3
    left = Inches(5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(1.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table.cell(0,0).text = "NT"
    table.cell(0,1).text = "FIRST"
    table.cell(0,2).text = "FOLLOW"
    
    data = [
        ["E", "id, (", "$, +, )"],
        ["T", "id, (", "$, +, *, )"],
        ["F", "id, (", "$, +, *, )"]
    ]
    
    for i in range(rows-1):
        table.cell(i+1, 0).text = data[i][0]
        table.cell(i+1, 1).text = data[i][1]
        table.cell(i+1, 2).text = data[i][2]

    # Style
    for cell in table.iter_cells():
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        if cell in [table.cell(0,0), table.cell(0,1), table.cell(0,2)]:
            cell.fill.fore_color.rgb = THEME_BLUE
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
            
    note = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    note.text_frame.text = "This grammar is SLR(1) because there are no conflicts in the parsing table."
    note.text_frame.paragraphs[0].font.size = Pt(16)
    note.text_frame.paragraphs[0].font.color.rgb = THEME_DARK_GRAY
    note.text_frame.paragraphs[0].font.italic = True
    
    add_notes(slide, "We use the standard Expression grammar. Compute FIRST and FOLLOW based on these rules.")

    # Slide 6: Solved Example - Canonical Collection (Simplified Graph)
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Solved Example: Canonical Collection (DFA)"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # We will draw a schematic representation of the DFA states
    # Drawing 9 states is crowded, so we group them by layers
    
    def draw_state_node(slide, id_val, x, y, items_preview):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.2), Inches(0.6))
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME_LIGHT_BLUE
        shape.line.color.rgb = THEME_BLUE
        tf = shape.text_frame
        tf.text = f"I{id_val}\n{items_preview}"
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[1].font.size = Pt(9)
        return shape

    # Layer 0
    s0 = draw_state_node(slide, 0, 1, 2, "E'->.E")
    
    # Layer 1 (on E)
    s1 = draw_state_node(slide, 1, 4, 2, "E'->E.")
    # Layer 1 (on T from 0)
    s3 = draw_state_node(slide, 3, 4, 4, "T->.T*F")
    
    # Layer 1 (on F from 3)
    s5 = draw_state_node(slide, 5, 4, 5.5, "F->.id")

    # Layer 2 (Shift +)
    s2 = draw_state_node(slide, 2, 6.5, 2, "E->E+.T")
    
    # Connectors
    slide.shapes.add_connector(MSO_SHAPE.ELLIPSE, s0.right, s0.top+0.2, s1.left, s1.top+0.2).line.color.rgb = THEME_DARK_GRAY
    slide.shapes.add_connector(MSO_SHAPE.ELLIPSE, s0.right, s0.bottom-0.1, s3.left, s3.top+0.2).line.color.rgb = THEME_DARK_GRAY
    
    # Additional transitions
    slide.shapes.add_connector(MSO_SHAPE.ELLIPSE, s1.right, s1.top+0.3, s2.left, s2.top+0.3).line.color.rgb = THEME_DARK_GRAY
    slide.shapes.add_connector(MSO_SHAPE.ELLIPSE, s3.right, s3.top+0.3, s5.left, s5.top+0.3).line.color.rgb = THEME_DARK_GRAY

    lbl = slide.shapes.add_textbox(Inches(7.5), Inches(2), Inches(2), Inches(2))
    lbl.text_frame.text = "Key States:\nI0: Start\nI1: Accept State\nI2, I3...: Intermediate"
    lbl.text_frame.paragraphs[0].font.size = Pt(14)

    add_notes(slide, "The DFA represents the LR(0) automaton. Nodes are sets of items. Transitions are on symbols.")

    # Slide 7: Solved Example - Parsing Table
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Solved Example: Parsing Table"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Create a representation of the table
    # Rows: State, Action (id, +, *, (, )), Goto (E, T, F)
    rows = 8
    cols = 9
    top = Inches(1.5)
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(5)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(0.5)
    
    # Headers
    headers = ["State", "id", "+", "*", "(", ")", "$", "E", "T", "F"]
    # Correction: cols is 9, headers is 10. Let's fix cols=10
    # Re-init table
    table = slide.shapes.add_table(rows, 10, left, top, width, height).table
    
    # Headers
    for i in range(10):
        cell = table.cell(0, i)
        cell.text = headers[i]
        cell.fill.solid()
        cell.fill.fore_color.rgb = THEME_BLUE
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Dummy Data for illustration
    data = [
        ["0", "s5", "", "", "s4", "", "", "1", "2", "3"],
        ["1", "", "s6", "", "", "", "acc", "", "", ""],
        ["2", "s5", "", "", "s4", "", "", "", "8", "3"],
        ["3", "", "r4", "s7", "", "r4", "r4", "", "", ""],
        ["4", "s5", "", "", "s4", "", "", "9", "2", "3"],
        ["5", "", "r6", "r6", "", "r6", "r6", "", "", ""],
        ["6", "s5", "", "", "s4", "", "", "", "", "10"]
    ]
    
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            table.cell(r_idx+1, c_idx).text = val
            table.cell(r_idx+1, c_idx).text_frame.paragraphs[0].font.size = Pt(10)
            table.cell(r_idx+1, c_idx).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Highlighting logic
            if "s" in val: table.cell(r_idx+1, c_idx).fill.fore_color.rgb = RGBColor(200, 230, 200) # Greenish
            if "r" in val: table.cell(r_idx+1, c_idx).fill.fore_color.rgb = RGBColor(230, 200, 200) # Reddish
            if "acc" in val: table.cell(r_idx+1, c_idx).fill.fore_color.rgb = RGBColor(255, 255, 150) # Yellow

    add_notes(slide, "sX = Shift to state X. rY = Reduce using production Y. acc = Accept.")

    # Slide 8: Solved Example - Parsing String
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Solved Example: Parsing \"id + id * id $\""
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Step-by-step Table
    rows, cols = 6, 4
    top = Inches(1.5)
    left = Inches(1)
    width = Inches(8)
    height = Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    headers = ["Stack", "Input", "Action", "Notes"]
    for i in range(4):
        table.cell(0,i).text = headers[i]
        table.cell(0,i).fill.fore_color.rgb = THEME_BLUE
        table.cell(0,i).text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
        table.cell(0,i).text_frame.paragraphs[0].font.size = Pt(14)

    # Steps (Simplified for slide space)
    steps_data = [
        ["0", "id + id * id $", "Shift", "Push id"],
        ["0 id 5", "+ id * id $", "Reduce F->id", "Pop 2, Push F(3)"],
        ["0 F 3", "+ id * id $", "Reduce T->F", "Pop 2, Push T(2)"],
        ["0 T 2", "+ id * id $", "Reduce E->T", "Pop 2, Push E(1)"],
        ["0 E 1", "+ id * id $", "Shift +", "Push +"]
    ]

    for r, row in enumerate(steps_data):
        for c, val in enumerate(row):
            table.cell(r+1, c).text = val
            table.cell(r+1, c).text_frame.paragraphs[0].font.size = Pt(12)

    add_notes(slide, "Walk through the parsing steps. Note how 'id' is reduced to F, then T, then E.")

    # Slide 9: Parsing Tree
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Solved Example: Parse Tree"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Draw Tree structure manually
    # E
    # / | \
    # E + T
    # |   / | \
    # T   T * F
    # |   |   |
    # F   F   id
    # |   |
    # id  id
    
    # Coordinates calculation
    y_root = 1.5
    y_l1 = 2.5
    y_l2 = 3.5
    y_l3 = 4.5
    y_l4 = 5.5
    
    # X positions
    x_center = 5
    x_left = 3
    x_right = 7
    x_l2_left = 2
    x_l2_mid = 4
    x_l2_right = 6
    x_l3_left = 2
    x_l3_right = 6
    x_l4_left = 2
    
    def add_tree_node(slide, text, x, y):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-0.4), Inches(y), Inches(0.8), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = THEME_LIGHT_BLUE
        shape.line.color.rgb = THEME_BLUE
        tf = shape.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        return shape

    # Nodes
    e1 = add_tree_node(slide, "E", x_center, y_root)
    e2 = add_tree_node(slide, "E", x_left, y_l1)
    plus = add_tree_node(slide, "+", x_center, y_l1)
    t1 = add_tree_node(slide, "T", x_right, y_l1)
    
    t2 = add_tree_node(slide, "T", x_left, y_l2)
    
    f1 = add_tree_node(slide, "F", x_left, y_l3)
    
    id1 = add_tree_node(slide, "id", x_left, y_l4)
    
    # Right subtree part
    t_sub = add_tree_node(slide, "T", x_l2_mid, y_l2)
    star = add_tree_node(slide, "*", x_right, y_l2)
    f_sub = add_tree_node(slide, "F", x_l2_right, y_l2)
    
    f_mid = add_tree_node(slide, "F", x_l2_mid, y_l3)
    id2 = add_tree_node(slide, "id", x_l2_mid, y_l4)
    
    id3 = add_tree_node(slide, "id", x_l2_right, y_l3)

    # Connectors (Lines)
    conns = [
        (e1, e2), (e1, plus), (e1, t1),
        (e2, t2),
        (t2, f1),
        (f1, id1),
        (t1, t_sub), (t1, star), (t1, f_sub),
        (t_sub, f_mid), (f_mid, id2),
        (f_sub, id3)
    ]
    
    for s, e in conns:
        # Use simple lines, adjusting centers roughly
        line = slide.shapes.add_shape(MSO_SHAPE.LINE, s.left+s.width/2, s.top+s.height, e.left+e.width/2, e.top)
        line.line.color.rgb = THEME_DARK_GRAY
        line.line.width = Pt(2)
        
    add_notes(slide, "Visual representation of the derivation E => E+T => T+T => F+T => id+T => id+T*F...")

    # Slide 10: Possible Conflicts
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Possible Conflicts in SLR(1)"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Shift-Reduce
    sr_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(2))
    sr_box.text_frame.text = "Shift-Reduce Conflict\n\nA state has:\n1. An item A -> α . β  (Shift on β)\n2. An item B -> γ . (Reduce)\n\nAND β is in FOLLOW(B)."
    sr_box.fill.solid()
    sr_box.fill.fore_color.rgb = RGBColor(255, 230, 230) # Light Red
    
    # Reduce-Reduce
    rr_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(2))
    rr_box.text_frame.text = "Reduce-Reduce Conflict\n\nA state has:\n1. An item A -> α .\n2. An item B -> β .\n\nAND FOLLOW(A) ∩ FOLLOW(B) ≠ Ø."
    rr_box.fill.solid()
    rr_box.fill.fore_color.rgb = RGBColor(230, 230, 255) # Light Blue

    # Diagram
    venn = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), Inches(4.5), Inches(2), Inches(2))
    venn2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.5), Inches(4.5), Inches(2), Inches(2))
    venn.fill.fore_color.rgb = RGBColor(200, 200, 255, 50)
    venn2.fill.fore_color.rgb = RGBColor(255, 200, 200, 50)
    
    lbl1 = slide.shapes.add_textbox(Inches(3.2), Inches(5.2), Inches(0.8), Inches(0.4))
    lbl1.text_frame.text = "FOLLOW(A)"
    lbl2 = slide.shapes.add_textbox(Inches(6), Inches(5.2), Inches(0.8), Inches(0.4))
    lbl2.text_frame.text = "FOLLOW(B)"
    
    conf_lbl = slide.shapes.add_textbox(Inches(4.6), Inches(4.5), Inches(1), Inches(0.4))
    conf_lbl.text_frame.text = "Conflict!"
    conf_lbl.text_frame.paragraphs[0].font.bold = True
    conf_lbl.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,0,0)

    add_notes(slide, "SLR(1) fails when the grammar is too ambiguous. Conflicts occur if the FOLLOW sets overlap.")

    # Slide 11: Shift-Reduce Conflict (Dangling Else)
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Example: Shift-Reduce Conflict (Dangling Else)"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Grammar
    g_text = "S -> if C then S | if C then S else S | a\nC -> b"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.6))
    tb.text_frame.text = g_text
    tb.text_frame.paragraphs[0].font.size = Pt(16)
    
    # Visualizing the Conflict
    # State with item S -> if C then S .  and S -> if C then S . else S
    # Actually the conflict is typically: after parsing "if C then S", we see "else".
    # We can Reduce S->if C then S, or Shift else.
    
    diag_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(2))
    diag_box.fill.solid()
    diag_box.fill.fore_color.rgb = RGBColor(245, 245, 245)
    diag_box.line.color.rgb = THEME_DARK_GRAY
    
    tf = diag_box.text_frame
    tf.text = "State I_k:\n1. S -> if C then S .    (Reduce on FOLLOW(S))\n2. S -> if C then S . else S  (Shift on 'else')"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[1].font.size = Pt(14)
    tf.paragraphs[1].font.color.rgb = RGBColor(200, 0, 0)
    tf.paragraphs[2].font.size = Pt(14)
    tf.paragraphs[2].font.color.rgb = RGBColor(0, 0, 200)
    
    # Two trees sketch
    t1 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(4), Inches(2))
    t1.text_frame.text = "Tree 1: Reduce\n(S matches the inner if)\nif C then [S] else a"
    
    t2 = slide.shapes.add_textbox(Inches(5), Inches(5), Inches(4), Inches(2))
    t2.text_frame.text = "Tree 2: Shift\n(S matches the outer if)\nif C then [if C then S else S]"
    
    add_notes(slide, "Dangling Else Problem. Shift vs Reduce on 'else'. Ambiguity arises because 'else' can belong to inner or outer 'if'.")

    # Slide 12: Reduce-Reduce Conflict
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Example: Reduce-Reduce Conflict"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    # Grammar
    g_text = "S -> A | B\nA -> x\nB -> x"
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.6))
    tb.text_frame.text = g_text
    tb.text_frame.paragraphs[0].font.size = Pt(16)

    # Diagram
    node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(2.5), Inches(4), Inches(1.5))
    node.fill.solid()
    node.fill.fore_color.rgb = RGBColor(255, 200, 200)
    node.line.color.rgb = RGBColor(255, 0, 0)
    node.text_frame.text = "State I_n:\nA -> x .\nB -> x ."
    node.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    node.text_frame.paragraphs[0].font.bold = True
    node.text_frame.paragraphs[0].font.size = Pt(16)

    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4), Inches(4.2), Inches(2), Inches(0.5))
    arrow.fill.fore_color.rgb = THEME_DARK_GRAY
    
    conflict_txt = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1))
    conflict_txt.text_frame.text = "If we see input 'x', should we reduce using A->x or B->x? Parser is confused."
    conflict_txt.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_notes(slide, "Ambiguous grammar where same terminal 'x' derives two different non-terminals.")

    # Slide 13: Resolving Conflicts
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Resolving Conflicts"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    left = Inches(0.5)
    top = Inches(1.2)
    
    strategies = [
        "1. Rewrite the Grammar:",
        "   - Make it unambiguous.",
        "   - Dangling Else Solution: Match each 'else' with closest 'if'.",
        "",
        "2. Use Stronger Parsers:",
        "   - CLR(1) / Canonical LR(1): Uses lookahead in items (more precise).",
        "   - LALR(1): Merges states of CLR(1) (Yacc/Bison standard).",
        "",
        "3. Manual Rules:",
        "   - Prefer Shift over Reduce (default in Yacc)."
    ]
    
    for i, text in enumerate(strategies):
        tb = slide.shapes.add_textbox(left, Inches(top + i*0.45), Inches(9), Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.name = FONT_NAME
        if text.startswith(str(i+1)[0]):
            p.font.bold = True
            p.font.color.rgb = THEME_BLUE
        else:
            p.font.color.rgb = THEME_DARK_GRAY

    add_notes(slide, "Rewriting is the best solution for correctness. Using CLR(1) is a technical fix but increases table size.")

    # Slide 14: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_box.text_frame.paragraphs[0].text = "Conclusion"
    title_box.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE
    title_box.text_frame.paragraphs[0].font.size = Pt(36)

    points = [
        "SLR(1) is a powerful parsing method.",
        "Construct table using LR(0) items + FOLLOW sets.",
        "Efficient for non-ambiguous grammars.",
        "Fails on conflicts (Shift-Reduce, Reduce-Reduce).",
        "LALR(1) (used in Bison) resolves many SLR(1) conflicts efficiently."
    ]
    
    for i, pt in enumerate(points):
        tb = slide.shapes.add_textbox(Inches(1), Inches(1.5 + i*0.6), Inches(8), Inches(0.6))
        tb.text_frame.text = pt
        tb.text_frame.paragraphs[0].font.size = Pt(20)
        tb.text_frame.paragraphs[0].font.name = FONT_NAME
        
    add_notes(slide, "Summary of the lecture. SLR(1) is the stepping stone to practical parser generators.")

    # Slide 15: Q&A
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Questions & Answers"
    slide.placeholders[1].text = "Thank you for your attention!"
    add_notes(slide, "Open floor for questions.")

    # Save
    file_path = "SLR1_Parsing_Compiler_Construction.pptx"
    prs.save(file_path)
    print(f"Presentation saved to: {file_path}")

if __name__ == "__main__":
    create_presentation()