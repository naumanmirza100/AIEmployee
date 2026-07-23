"""
Document Processing Agent for Operations Agent
Handles document upload, parsing, text extraction, chunking, entity extraction,
and auto-classification for PDF, DOCX, Excel, PPT, CSV, and TXT files.
"""

import os
import hashlib
import logging
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from django.conf import settings
from django.utils import timezone

from marketing_agent.agents.marketing_base_agent import MarketingBaseAgent

logger = logging.getLogger(__name__)

# Chunk config
CHUNK_SIZE = 1000       # ~1000 chars per chunk
CHUNK_OVERLAP = 150     # overlap between chunks


def _invalidate_operations_indexes(company_id, has_embeddings: bool) -> None:
    """After a doc goes live, rebuild the FAISS index and clear the answer cache
    for its company. Both are best-effort — a cache miss must never break an
    otherwise successful upload."""
    if not company_id:
        return
    if has_embeddings:
        try:
            from operations_agent.vector_store import mark_index_dirty
            mark_index_dirty(company_id)
        except Exception:
            logger.exception("Operations: failed to mark FAISS index dirty")
    try:
        from operations_agent.agents.knowledge_qa_agent import invalidate_answer_cache_for_company
        invalidate_answer_cache_for_company(company_id)
    except Exception:
        logger.exception("Operations: failed to invalidate answer cache")


class DocumentProcessingAgent(MarketingBaseAgent):
    """
    Processes uploaded documents: extracts text, metadata, entities,
    auto-classifies type, chunks for RAG, and optionally generates embeddings.
    """

    SUPPORTED_EXTENSIONS = {
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.xlsx': 'xlsx',
        '.csv': 'csv',
        '.pptx': 'pptx',
        '.txt': 'txt',
        '.md': 'txt',
    }

    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50 MB

    def __init__(self):
        super().__init__(use_embeddings=False)
        self.agent_name = 'DocumentProcessingAgent'

    # ------------------------------------------------------------------
    # Main entry
    # ------------------------------------------------------------------
    def process(self, action: str, **kwargs) -> Dict:
        actions = {
            'process_file': self._process_file,
            'classify': self._classify_document,
            'extract_entities': self._extract_entities,
        }
        handler = actions.get(action)
        if not handler:
            return {'success': False, 'error': f'Unknown action: {action}'}
        try:
            return handler(**kwargs)
        except Exception as e:
            from core.api_key_service import KeyServiceError
            if isinstance(e, KeyServiceError):
                raise
            logger.error(f'{self.agent_name} error: {e}', exc_info=True)
            return {'success': False, 'error': str(e)}

    # ------------------------------------------------------------------
    # File processing pipeline
    # ------------------------------------------------------------------
    def _process_file(self, file_path: str, original_filename: str,
                      company_id: int, uploaded_by_id: int,
                      title: str = '', tags: str = '',
                      existing_doc_id: int = None) -> Dict:
        """Full pipeline: validate → extract → classify → entities → chunk → save.

        When ``existing_doc_id`` is given (async upload flow), that placeholder
        row is stamped ``processing`` and updated in place instead of creating a
        new one; failures stamp it ``failed`` so the status poll can report them.
        """
        from operations_agent.models import OperationsDocument, OperationsDocumentChunk
        from core.models import Company, CompanyUser

        def _fail(doc_row, message):
            """Stamp a placeholder row failed (async flow) and return an error."""
            if doc_row is not None:
                try:
                    doc_row.processing_status = 'failed'
                    doc_row.processing_error = str(message)[:2000]
                    doc_row.save(update_fields=['processing_status', 'processing_error', 'updated_at'])
                except Exception:
                    logger.exception("Operations: failed to stamp doc %s as failed", getattr(doc_row, 'id', '?'))
            return {'success': False, 'error': message}

        placeholder = None
        if existing_doc_id:
            placeholder = OperationsDocument.objects.filter(id=existing_doc_id).first()
            if placeholder:
                placeholder.processing_status = 'processing'
                placeholder.save(update_fields=['processing_status', 'updated_at'])

        # 1. Validate
        ext = Path(original_filename).suffix.lower()
        file_type = self.SUPPORTED_EXTENSIONS.get(ext)
        if not file_type:
            return _fail(placeholder, f'Unsupported file type: {ext}. Supported: {", ".join(self.SUPPORTED_EXTENSIONS.keys())}')

        file_size = os.path.getsize(file_path)
        if file_size > self.MAX_FILE_SIZE:
            return _fail(placeholder, f'File too large ({file_size / 1024 / 1024:.1f} MB). Max: 50 MB')

        # 2. Extract text
        success, extracted_text, page_count, error = self._extract_text(file_path, file_type)
        if not success:
            return _fail(placeholder, error)

        if not extracted_text or not extracted_text.strip():
            return _fail(placeholder, 'No text could be extracted from this document')

        # 3. Auto-classify document type via LLM
        doc_type = self._classify_document(text=extracted_text[:3000])

        # 4. Extract entities via LLM
        entities = self._extract_entities(text=extracted_text[:4000])

        # 4b. Generate summary & key insights via LLM
        summary_result = self._generate_summary(text=extracted_text[:5000])
        summary = summary_result.get('summary', '') if isinstance(summary_result, dict) else ''
        key_insights = summary_result.get('key_insights', []) if isinstance(summary_result, dict) else []

        # 5. Build metadata
        metadata = {
            'original_filename': original_filename,
            'file_extension': ext,
            'char_count': len(extracted_text),
            'word_count': len(extracted_text.split()),
            'page_count': page_count,
        }

        # 6. Save OperationsDocument
        company = Company.objects.get(pk=company_id)
        uploaded_by = CompanyUser.objects.get(pk=uploaded_by_id) if uploaded_by_id else None

        resolved_doc_type = doc_type.get('document_type', 'other') if isinstance(doc_type, dict) else 'other'
        doc_fields = dict(
            company=company,
            uploaded_by=uploaded_by,
            title=title or Path(original_filename).stem,
            original_filename=original_filename,
            file=file_path,
            file_type=file_type,
            document_type=resolved_doc_type,
            file_size=file_size,
            page_count=page_count,
            parsed_text=extracted_text,
            summary=summary,
            key_insights=key_insights,
            metadata=metadata,
            entities=entities.get('entities', {}) if isinstance(entities, dict) else {},
            tags=tags,
            is_processed=True,
            processed_at=timezone.now(),
            processing_status='processing',
        )
        if placeholder is not None:
            # Async flow: fill in the pre-created placeholder row.
            for k, v in doc_fields.items():
                setattr(placeholder, k, v)
            placeholder.save()
            doc = placeholder
        else:
            doc = OperationsDocument.objects.create(**doc_fields)

        # 7. Chunk text (section-aware, TOC/junk filtered) + embed + save
        chunk_count, embedded, embed_model = self._chunk_embed_and_store(
            doc, extracted_text, page_count, resolved_doc_type,
        )

        # 8. Drop the source file — everything we need (parsed_text, chunks,
        #    summary, insights) is already in the DB, and reindex works from
        #    parsed_text. Keeping the binary just wastes disk.
        try:
            src = Path(file_path)
            if src.exists():
                src.unlink()
        except OSError:
            logger.warning("Operations: could not delete source file %s", file_path)

        # 9. Stamp final RAG state + invalidate FAISS / answer cache
        doc.chunks_total = chunk_count
        doc.chunks_processed = chunk_count
        doc.is_indexed = embedded
        doc.embedding_model = embed_model or ''
        doc.processing_status = 'ready'
        doc.file = ''  # file removed above — don't advertise a dead path
        doc.save(update_fields=[
            'chunks_total', 'chunks_processed', 'is_indexed',
            'embedding_model', 'processing_status', 'file', 'updated_at',
        ])
        _invalidate_operations_indexes(doc.company_id, embedded)

        self.log_action('process_file', {
            'document_id': doc.id,
            'filename': original_filename,
            'chunks': chunk_count,
            'embedded': embedded,
        })

        return {
            'success': True,
            'document': {
                'id': doc.id,
                'title': doc.title,
                'file_type': doc.file_type,
                'document_type': doc.document_type,
                'page_count': page_count,
                'word_count': metadata['word_count'],
                'chunks_created': chunk_count,
                'entities': doc.entities,
                'summary': summary,
                'key_insights': key_insights,
            },
        }

    # ------------------------------------------------------------------
    # Text extraction
    # ------------------------------------------------------------------
    def _extract_text(self, file_path: str, file_type: str) -> Tuple[bool, str, int, Optional[str]]:
        """Returns (success, text, page_count, error)."""
        try:
            if file_type == 'pdf':
                return self._extract_pdf(file_path)
            elif file_type == 'docx':
                return self._extract_docx(file_path)
            elif file_type == 'xlsx':
                return self._extract_xlsx(file_path)
            elif file_type == 'csv':
                return self._extract_csv(file_path)
            elif file_type == 'pptx':
                return self._extract_pptx(file_path)
            elif file_type == 'txt':
                return self._extract_txt(file_path)
            else:
                return False, '', 0, f'Unsupported file type: {file_type}'
        except Exception as e:
            logger.error(f'Text extraction failed for {file_path}: {e}', exc_info=True)
            return False, '', 0, str(e)

    @staticmethod
    def _extract_pdf(file_path: str) -> Tuple[bool, str, int, Optional[str]]:
        try:
            import pdfplumber
            pages_text = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ''
                    pages_text.append(text)
            return True, '\n\n'.join(pages_text), len(pages_text), None
        except ImportError:
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                pages_text = [p.extract_text() or '' for p in reader.pages]
                return True, '\n\n'.join(pages_text), len(pages_text), None
            except ImportError:
                return False, '', 0, 'No PDF library available (install pdfplumber or PyPDF2)'
        except Exception as e:
            return False, '', 0, f'PDF extraction failed: {e}'

    @staticmethod
    def _extract_docx(file_path: str) -> Tuple[bool, str, int, Optional[str]]:
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            page_count = max(1, len(paragraphs) // 40)  # ~40 paragraphs per page estimate
            return True, '\n\n'.join(paragraphs), page_count, None
        except ImportError:
            return False, '', 0, 'python-docx not installed'
        except Exception as e:
            return False, '', 0, f'DOCX extraction failed: {e}'

    @staticmethod
    def _extract_xlsx(file_path: str) -> Tuple[bool, str, int, Optional[str]]:
        try:
            import pandas as pd
            xls = pd.ExcelFile(file_path)
            all_text = []
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                all_text.append(f'--- Sheet: {sheet_name} ---')
                all_text.append(df.to_string(index=False))
            return True, '\n\n'.join(all_text), len(xls.sheet_names), None
        except ImportError:
            return False, '', 0, 'pandas/openpyxl not installed'
        except Exception as e:
            return False, '', 0, f'Excel extraction failed: {e}'

    @staticmethod
    def _extract_csv(file_path: str) -> Tuple[bool, str, int, Optional[str]]:
        try:
            import pandas as pd
            df = pd.read_csv(file_path)
            return True, df.to_string(index=False), 1, None
        except ImportError:
            return False, '', 0, 'pandas not installed'
        except Exception as e:
            return False, '', 0, f'CSV extraction failed: {e}'

    @staticmethod
    def _extract_pptx(file_path: str) -> Tuple[bool, str, int, Optional[str]]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        texts.append(shape.text)
                if texts:
                    slides_text.append(f'--- Slide {i} ---\n' + '\n'.join(texts))
            return True, '\n\n'.join(slides_text), len(prs.slides), None
        except ImportError:
            return False, '', 0, 'python-pptx not installed'
        except Exception as e:
            return False, '', 0, f'PPTX extraction failed: {e}'

    @staticmethod
    def _extract_txt(file_path: str) -> Tuple[bool, str, int, Optional[str]]:
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            page_count = max(1, len(text) // 3000)
            return True, text, page_count, None
        except Exception as e:
            return False, '', 0, f'Text extraction failed: {e}'

    # ------------------------------------------------------------------
    # Chunking + embedding (RAG)
    # ------------------------------------------------------------------
    def _chunk_embed_and_store(self, doc, text: str, page_count: int,
                               document_type: str):
        """Chunk (section-aware, TOC filtered), embed in batches, and persist.

        Returns ``(chunk_count, embedded_bool, embedding_model)``. Falls back to
        no-embedding storage when the embedding provider is unavailable, so the
        keyword retrieval path still works. Mirrors the HR/Frontline pipeline.
        """
        from operations_agent.models import OperationsDocumentChunk

        chunk_pairs = self._chunk_pairs(text, page_count, document_type)
        if not chunk_pairs:
            return 0, False, ''

        # Embedding is optional — the retriever falls back to keyword search
        # when a chunk has no vector. Provider is env-keyed and shared across
        # HR / Frontline / Operations.
        embedding_service = None
        has_embeddings = False
        embed_model = ''
        try:
            from core.Frontline_agent.embedding_service import EmbeddingService
            embedding_service = EmbeddingService()
            has_embeddings = embedding_service.is_available()
            embed_model = getattr(embedding_service, 'embedding_model', '') or ''
        except Exception:
            logger.warning("Operations: embedding service unavailable; storing chunks without vectors")

        texts = [c for c, _h, _p in chunk_pairs]
        embeddings = [None] * len(texts)
        if has_embeddings:
            batch_size = 20
            for i in range(0, len(texts), batch_size):
                batch = texts[i:i + batch_size]
                try:
                    batch_vecs = embedding_service.generate_embeddings_batch(batch)
                except Exception as exc:
                    logger.warning("Operations embedding batch failed: %s", exc)
                    batch_vecs = [None] * len(batch)
                for j, v in enumerate(batch_vecs or []):
                    embeddings[i + j] = v

        rows = []
        any_vector = False
        for i, (content, heading, page) in enumerate(chunk_pairs):
            emb = embeddings[i]
            if emb:
                any_vector = True
            rows.append(OperationsDocumentChunk(
                document=doc,
                chunk_index=i,
                content=content,
                section_heading=(heading or '')[:300],
                page_number=page,
                token_count=len(content.split()),
                embedding=emb,
            ))
        OperationsDocumentChunk.objects.bulk_create(rows)
        return len(rows), any_vector, (embed_model if any_vector else '')

    def _chunk_pairs(self, text: str, page_count: int, document_type: str):
        """Return a list of ``(content, heading, page_number)`` tuples.

        Section-aware doc types use the heading-aware chunker (which also drops
        TOC/index junk); everything else uses a fixed-window chunker that still
        skips junk chunks. Page numbers are estimated proportionally.
        """
        from operations_agent.chunking import (
            chunk_with_headings, looks_like_toc_or_index, SECTION_AWARE_TYPES,
        )

        total = max(1, len(text))

        def _page_for(offset: int) -> int:
            return min(max(page_count, 1), max(1, int((offset / total) * max(page_count, 1)) + 1))

        pairs = []
        if document_type in SECTION_AWARE_TYPES:
            for content, heading in chunk_with_headings(
                text, max_chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP,
            ):
                # Recover an approximate page from where the chunk starts.
                off = text.find(content[:60]) if content else -1
                page = _page_for(off if off >= 0 else 0)
                pairs.append((content, heading, page))
            if pairs:
                return pairs

        # Fixed-window fallback (also used when section chunking yielded nothing)
        idx_start = 0
        step = max(1, CHUNK_SIZE - CHUNK_OVERLAP)
        while idx_start < len(text):
            content = text[idx_start:idx_start + CHUNK_SIZE].strip()
            if content and not looks_like_toc_or_index(content):
                pairs.append((content, '', _page_for(idx_start)))
            idx_start += step
        return pairs

    # ------------------------------------------------------------------
    # LLM-based classification
    # ------------------------------------------------------------------
    def _classify_document(self, text: str = '', **kwargs) -> Dict:
        """Classify document type using LLM."""
        if not text:
            return {'document_type': 'other'}
        try:
            prompt = f"""Classify the following document into exactly ONE of these types:
report, invoice, contract, memo, spreadsheet, presentation, policy, manual, other

Document excerpt:
\"\"\"
{text[:2000]}
\"\"\"

Respond with ONLY the type name (one word, lowercase). Nothing else."""

            result = self._call_llm_for_reasoning(prompt, temperature=0.1, max_tokens=20)
            doc_type = result.strip().lower().replace('.', '').replace('"', '').split()[0] if result else 'other'
            valid_types = ['report', 'invoice', 'contract', 'memo', 'spreadsheet', 'presentation', 'policy', 'manual', 'other']
            if doc_type not in valid_types:
                doc_type = 'other'
            return {'document_type': doc_type}
        except Exception as e:
            from core.api_key_service import KeyServiceError
            if isinstance(e, KeyServiceError):
                raise
            logger.warning(f'Classification failed, defaulting to other: {e}')
            return {'document_type': 'other'}

    # ------------------------------------------------------------------
    # LLM-based entity extraction
    # ------------------------------------------------------------------
    def _extract_entities(self, text: str = '', **kwargs) -> Dict:
        """Extract key entities from document text using LLM."""
        if not text:
            return {'entities': {}}
        try:
            prompt = f"""Extract key entities from the following document. Return a JSON object with these keys:
- dates: list of important dates found
- amounts: list of monetary amounts found
- names: list of person names found
- organizations: list of organization/company names found
- key_terms: list of important domain-specific terms

Document excerpt:
\"\"\"
{text[:3000]}
\"\"\"

Return ONLY valid JSON. No explanation."""

            result = self._call_llm_for_reasoning(prompt, temperature=0.1, max_tokens=500)

            import json
            # Try to parse JSON from response
            cleaned = result.strip()
            if cleaned.startswith('```'):
                cleaned = cleaned.split('\n', 1)[-1].rsplit('```', 1)[0]
            entities = json.loads(cleaned)
            return {'entities': entities}
        except Exception as e:
            from core.api_key_service import KeyServiceError
            if isinstance(e, KeyServiceError):
                raise
            logger.warning(f'Entity extraction failed: {e}')
            return {'entities': {}}

    # ------------------------------------------------------------------
    # LLM-based summary generation
    # ------------------------------------------------------------------
    def _generate_summary(self, text: str = '', **kwargs) -> Dict:
        """Generate a concise summary and key insights from document text."""
        if not text:
            return {'summary': '', 'key_insights': []}
        try:
            prompt = f"""Analyze the following document and provide:
1. A clear, concise summary (3-5 sentences) that captures the main purpose, content, and conclusions.
2. A list of 3-6 key insights or findings from the document.

Document excerpt:
\"\"\"
{text[:4500]}
\"\"\"

Return ONLY valid JSON in this exact format:
{{
  "summary": "Your summary here...",
  "key_insights": ["Insight 1", "Insight 2", "Insight 3"]
}}

No explanation outside the JSON."""

            result = self._call_llm_for_reasoning(prompt, temperature=0.2, max_tokens=800)

            import json
            cleaned = result.strip()
            if cleaned.startswith('```'):
                cleaned = cleaned.split('\n', 1)[-1].rsplit('```', 1)[0]
            parsed = json.loads(cleaned)
            return {
                'summary': parsed.get('summary', ''),
                'key_insights': parsed.get('key_insights', []),
            }
        except Exception as e:
            from core.api_key_service import KeyServiceError
            if isinstance(e, KeyServiceError):
                raise
            logger.warning(f'Summary generation failed: {e}')
            return {'summary': '', 'key_insights': []}
